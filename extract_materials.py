"""Extract inventory, Excel formulas, and document summaries from BASE folder."""
import json
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

BASE = Path(r"C:\Users\alexd\Projects\formulacion\BASE")
OUT = Path(r"C:\Users\alexd\Projects\formulacion\materials_report.json")


def inventory():
    files = []
    for f in sorted(BASE.iterdir()):
        if f.is_file():
            files.append({
                "name": f.name,
                "path": str(f),
                "size_bytes": f.stat().st_size,
                "ext": f.suffix.lower(),
            })
    return files


def extract_xlsx_formulas(path: Path):
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=False)
    formulas = []
    for sn in wb.sheetnames:
        ws = wb[sn]
        for row in ws.iter_rows():
            for cell in row:
                val = cell.value
                if val is None:
                    continue
                is_formula = cell.data_type == "f" or (
                    isinstance(val, str) and val.startswith("=")
                )
                if is_formula:
                    formulas.append({
                        "sheet": sn,
                        "cell": cell.coordinate,
                        "formula": str(val),
                    })
    return formulas


def extract_xls_formulas(path: Path):
    import xlrd

    wb = xlrd.open_workbook(str(path), formatting_info=False)
    formulas = []
    for si in range(wb.nsheets):
        sh = wb.sheet_by_index(si)
        for r in range(sh.nrows):
            for c in range(sh.ncols):
                cell = sh.cell(r, c)
                val = cell.value
                if isinstance(val, str) and val.startswith("="):
                    formulas.append({
                        "sheet": sh.name,
                        "cell": f"{xlrd.colname(c)}{r + 1}",
                        "formula": val,
                    })
    return formulas


def skim_pptx(path: Path, max_slides=20):
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides[:max_slides], 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append({"slide": i, "text": texts[:5]})
    return {"slide_count": len(prs.slides), "sample_slides": slides}


def skim_docx(path: Path, max_paras=30):
    from docx import Document

    doc = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return {"paragraph_count": len(paras), "sample_paragraphs": paras[:max_paras]}


def skim_xlsx_sheets(path: Path):
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = {}
    for sn in wb.sheetnames:
        ws = wb[sn]
        headers = []
        for row in ws.iter_rows(min_row=1, max_row=5, values_only=True):
            headers.append([str(c) if c is not None else "" for c in row[:12]])
        sheets[sn] = {"preview_rows": headers}
    return sheets


def main():
    report = {"inventory": inventory(), "excel": {}, "documents": {}}

    for f in BASE.iterdir():
        if not f.is_file():
            continue
        ext = f.suffix.lower()
        name = f.name

        if ext in {".xlsx", ".xlsm"}:
            try:
                formulas = extract_xlsx_formulas(f)
                unique = {}
                for item in formulas:
                    unique.setdefault(item["formula"], []).append(
                        f"{item['sheet']}!{item['cell']}"
                    )
                report["excel"][name] = {
                    "formula_count": len(formulas),
                    "unique_formula_count": len(unique),
                    "formulas": [
                        {"formula": k, "locations": v[:5], "count": len(v)}
                        for k, v in sorted(unique.items(), key=lambda x: -len(x[1]))
                    ],
                    "sheets_preview": skim_xlsx_sheets(f),
                }
            except Exception as e:
                report["excel"][name] = {"error": str(e)}

        elif ext == ".xls":
            try:
                formulas = extract_xls_formulas(f)
                unique = {}
                for item in formulas:
                    unique.setdefault(item["formula"], []).append(
                        f"{item['sheet']}!{item['cell']}"
                    )
                report["excel"][name] = {
                    "formula_count": len(formulas),
                    "unique_formula_count": len(unique),
                    "formulas": [
                        {"formula": k, "locations": v[:5], "count": len(v)}
                        for k, v in sorted(unique.items(), key=lambda x: -len(x[1]))
                    ],
                }
            except Exception as e:
                report["excel"][name] = {"error": str(e)}

        elif ext == ".pptx":
            try:
                report["documents"][name] = skim_pptx(f)
            except Exception as e:
                report["documents"][name] = {"error": str(e)}

        elif ext == ".docx":
            try:
                report["documents"][name] = skim_docx(f)
            except Exception as e:
                report["documents"][name] = {"error": str(e)}

    OUT.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Report written to {OUT}")
    print(f"Files: {len(report['inventory'])}")
    print(f"Excel files: {len(report['excel'])}")
    print(f"Documents: {len(report['documents'])}")


if __name__ == "__main__":
    main()
